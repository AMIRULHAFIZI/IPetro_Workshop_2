import os
import re
import json
import click
import uuid 
import time
import copy
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
from flask import Flask, request, render_template, send_from_directory, send_file, redirect, url_for, flash
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from flask_migrate import Migrate
from datetime import datetime, time
from sqlalchemy import func
from dotenv import load_dotenv
from pptx.util import Pt, Inches  # <--- Make sure Inches is added here
from pptx.enum.text import PP_ALIGN

import pandas as pd
import pdf2image
from PIL import Image
import requests
import base64
from openpyxl import load_workbook
from openpyxl.styles import Alignment 
import shutil
import time as sleep_time 
import tempfile
from pptx.util import Inches

from flask_wtf import FlaskForm
from wtforms import StringField, PasswordField, SubmitField, SelectField, TextAreaField, BooleanField
from wtforms.validators import DataRequired, Email, EqualTo, Length, ValidationError, Optional
from email_validator import validate_email, EmailNotValidError
from flask_wtf.file import FileField, FileAllowed
from werkzeug.utils import secure_filename
from functools import wraps
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from flask_mail import Mail, Message
from PIL import Image

# --- LOAD ENVIRONMENT VARIABLES (Reading from .env) ---
load_dotenv()
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
DATABASE_URL = os.environ.get('DATABASE_URL')
POPPLER_PATH = os.environ.get('POPPLER_PATH')

if not GEMINI_API_KEY:
    raise ValueError("CRITICAL ERROR: 'GEMINI_API_KEY' not set in .env file.")
if not DATABASE_URL:
    raise ValueError("CRITICAL ERROR: 'DATABASE_URL' not set in .env file.")
if not POPPLER_PATH:
    raise ValueError("CRITICAL ERROR: 'POPPLER_PATH' not set in .env file.")

# --- APP SETUP ---
app = Flask(__name__)
app.config['SECRET_KEY'] = 'a-super-secret-key-that-is-hard-to-guess'
app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

# --- DB CONNECTION FIX ---
app.config['SQLALCHEMY_ENGINE_OPTIONS'] = {
    'pool_size': 10,       # Keep up to 10 connections in the pool
    'pool_recycle': 280,   # Recycle connections before the default 300s timeout
    'pool_pre_ping': True, # Check if the DB is alive before trying to use it
}

# --- MAIL CONFIGURATION ---
app.config['MAIL_SERVER'] = 'smtp.gmail.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USERNAME'] = os.environ.get('MAIL_USERNAME')
app.config['MAIL_PASSWORD'] = os.environ.get('MAIL_PASSWORD')
app.config['MAIL_DEFAULT_SENDER'] = os.environ.get('MAIL_USERNAME')

mail = Mail(app)
db = SQLAlchemy(app)
migrate = Migrate(app, db)

login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = 'Please log in to access this page.'
login_manager.login_message_category = 'error'

# --- FOLDER CONFIG ---
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ANNOUNCEMENT_FOLDER = 'announcements'

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(ANNOUNCEMENT_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['ANNOUNCEMENT_FOLDER'] = ANNOUNCEMENT_FOLDER

# --- EXCEL TEMPLATE CONFIG ---
TEMPLATE_FILE = 'EXCEL.xlsx'
TARGET_SHEET = 'Sheet1'
template_column_order = [
    'NO.', 'EQUIPMENT NO. ', 'PMT NO.', 'EQUIPMENT DESCRIPTION', 'PARTS',
    'PHASE', 'FLUID', 'TYPE', 'SPEC.', 'GRADE',
    'INSULATION',
    'TEMP. (°C) ', 'PRESSURE (Mpa) ', 
    'TEMP. (°C)', 'PRESSURE (Mpa)' 
]

# ==================== HELPER FUNCTIONS ==================

def _set_cell_text(cell, text, font_size=10):
    """
    Helper to set text in a table cell with specific font size and alignment.
    """
    if text is None:
        text = "-"
    
    text_str = str(text).strip()
    
    if text_str == '' or text_str.lower() == 'nan':
        text_str = "-"
        
    cell.text_frame.text = text_str
    
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.name = 'Arial'
        
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_new_row(table):
    """
    Copies the last row of the table and appends it to the end.
    """
    import copy
    new_row_xml = copy.deepcopy(table._tbl.tr_lst[-1])
    table._tbl.append(new_row_xml)
    new_row = table.rows[len(table.rows) - 1]
    for cell in new_row.cells:
        cell.text_frame.text = ""
    return new_row

def duplicate_slide(pres, index):
    """
    Robustly duplicate the slide at 'index', handling images correctly.
    """
    source = pres.slides[index]
    dest = pres.slides.add_slide(source.slide_layout)

    for shp in dest.shapes:
        shp.element.getparent().remove(shp.element)

    for shp in source.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shp.image.blob
                image_stream = io.BytesIO(blob)
                dest.shapes.add_picture(
                    image_stream, shp.left, shp.top, shp.width, shp.height
                )
            except Exception as e:
                print(f"Warning: Could not copy image on slide duplication: {e}")
        else:
            new_el = copy.deepcopy(shp.element)
            dest.shapes._spTree.append(new_el)

    return dest

def generate_ppt_internal(filename):
    """
    Internal function to generate the PPT file with interleaved slides:
    [Data Slides A] -> [Drawing A] -> [Data Slides B] -> [Drawing B]
    """
    # 1. Path Setup
    excel_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    template_path = os.path.join(app.root_path, 'Inspection Plan Powerpoint Template.pptx')
    output_pptx_name = f"Inspection_Plan_{filename.replace('.xlsx', '.pptx')}"
    output_pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], output_pptx_name)

    # 2. Check Files
    if not os.path.exists(excel_path) or not os.path.exists(template_path):
        return None

    # ==========================================
    # DATA PROCESSING
    # ==========================================
    try:
        df_raw = pd.read_excel(excel_path, header=None)
        
        # 1. Find the Anchor Row
        anchor_idx = None
        for idx, row in df_raw.iterrows():
            r_str = row.astype(str).str.upper().tolist()
            if "TYPE" in r_str and "SPEC." in r_str:
                anchor_idx = idx
                break
        
        if anchor_idx is None: return None

        # 2. Slice Data
        df_data = df_raw.iloc[anchor_idx + 1:].copy().reset_index(drop=True)
        
        header_row_sub = df_raw.iloc[anchor_idx].astype(str).str.upper().tolist()
        header_row_main = []
        if anchor_idx > 0:
            header_row_main = df_raw.iloc[anchor_idx-1].astype(str).str.upper().tolist()

        # 3. IDENTIFY COLUMNS
        idx_fluid, idx_part = None, None
        idx_type, idx_spec, idx_grade = None, None, None
        idx_insul, idx_temp, idx_press = None, None, None
        idx_equip_no, idx_pmt_no, idx_desc = None, None, None

        num_cols = df_data.shape[1]
        
        for c in range(num_cols):
            h_main = header_row_main[c] if c < len(header_row_main) else ""
            h_sub = header_row_sub[c] if c < len(header_row_sub) else ""
            full_head = (h_main + " " + h_sub).upper()

            if "EQUIPMENT NO" in full_head or "EQUIP NO" in full_head: idx_equip_no = c
            if "PMT" in full_head: idx_pmt_no = c
            if "DESCRIPTION" in full_head or "DESC" in full_head: idx_desc = c
            if "FLUID" in full_head: idx_fluid = c
            if "INSUL" in full_head: idx_insul = c
            if "SPEC" in full_head: idx_spec = c
            if "GRADE" in full_head or "GR." in full_head or h_sub == "GR": idx_grade = c
            if "TEMP" in full_head and "DESIGN" not in full_head: idx_temp = c
            if "PRESS" in full_head and "DESIGN" not in full_head: idx_press = c
            if "TYPE" in h_sub: idx_type = c
            if "PARTS" in full_head: idx_part = c

        # Fallbacks
        if idx_equip_no is None: idx_equip_no = 1
        if idx_pmt_no is None: idx_pmt_no = 2
        if idx_desc is None: idx_desc = 3
        if idx_part is None: idx_part = 4
        if idx_fluid is None: idx_fluid = 6

        # 4. FIX: FORWARD FILL METADATA (Crucial for Grouping)
        cols_to_fill = []
        if idx_equip_no is not None: cols_to_fill.append(idx_equip_no)
        if idx_pmt_no is not None: cols_to_fill.append(idx_pmt_no)
        if idx_desc is not None: cols_to_fill.append(idx_desc)

        for col in cols_to_fill:
             if col < df_data.shape[1]:
                df_data.iloc[:, col] = df_data.iloc[:, col].replace(['', '-', 'nan', 'None'], pd.NA)
                df_data.iloc[:, col] = df_data.iloc[:, col].ffill()
                df_data.iloc[:, col] = df_data.iloc[:, col].fillna('')

        df_data = df_data.fillna('-')

        # Filter Empty Rows
        if idx_part is not None:
             df_data = df_data[
                (df_data[idx_part].astype(str).str.strip() != '-') & 
                (df_data[idx_part].astype(str).str.strip() != '') 
            ]

    except Exception as e:
        print(f"PPT Generation Error (Data): {e}")
        return None

    # ==========================================
    # GENERATE POWERPOINT
    # ==========================================
    try:
        prs = Presentation(template_path)
        MAX_ROWS = 5
        
        # --- NEW STRATEGY: GROUP BY PMT NO (FILE) ---
        # This ensures we process [Data -> Drawing] for File A, then [Data -> Drawing] for File B
        if idx_pmt_no is not None:
            # Get the column name to group by
            pmt_col = df_data.columns[idx_pmt_no]
            # Group by PMT No, keeping the order (sort=False)
            grouped = df_data.groupby(pmt_col, sort=False)
        else:
            # Fallback if PMT column fails: treat everything as one group
            grouped = [("Unknown", df_data)]

        for pmt_val, group_df in grouped:
            
            # --- 1. GENERATE DATA SLIDES FOR THIS GROUP ---
            data_rows = [row for _, row in group_df.iterrows()]
            chunks = [data_rows[i:i + MAX_ROWS] for i in range(0, len(data_rows), MAX_ROWS)] if data_rows else [[]]
            
            # Capture metadata from the first row of this group
            first_row = data_rows[0]
            current_tag = str(first_row.get(idx_equip_no, '')).strip()
            current_pmt = str(first_row.get(idx_pmt_no, '')).strip()
            current_desc = str(first_row.get(idx_desc, '')).strip()
            
            # Determine Drawing Filename
            current_drawing_filename = ""
            if current_pmt and current_pmt != '-':
                current_drawing_filename = current_pmt + ".pdf"

            for chunk in chunks:
                # Copy Master Slide (Index 0)
                slide = duplicate_slide(prs, 0)
                
                # Overlay Text Boxes
                txBox_desc = slide.shapes.add_textbox(Inches(2.6), Inches(0.55), Inches(3.5), Inches(0.4))
                tf_desc = txBox_desc.text_frame
                tf_desc.text = current_desc
                tf_desc.paragraphs[0].font.size = Pt(10)
                tf_desc.paragraphs[0].font.name = 'Arial'
                tf_desc.paragraphs[0].font.bold = True

                txBox_tag = slide.shapes.add_textbox(Inches(6.3), Inches(0.55), Inches(1.8), Inches(0.4))
                tf_tag = txBox_tag.text_frame
                tf_tag.text = current_tag
                tf_tag.paragraphs[0].font.size = Pt(11) 
                tf_tag.paragraphs[0].font.name = 'Arial'
                tf_tag.paragraphs[0].font.bold = True

                txBox_pmt = slide.shapes.add_textbox(Inches(9.5), Inches(0.55), Inches(1.8), Inches(0.4))
                tf_pmt = txBox_pmt.text_frame
                tf_pmt.text = current_pmt
                tf_pmt.paragraphs[0].font.size = Pt(10)
                tf_pmt.paragraphs[0].font.name = 'Arial'
                tf_pmt.paragraphs[0].font.bold = True

                # Fill Table
                main_table = None
                for shape in slide.shapes:
                    if shape.has_table:
                        try:
                            r0 = " ".join([c.text_frame.text.upper() for c in shape.table.rows[0].cells])
                            if "COMPONENT" in r0 and "FLUID" in r0:
                                main_table = shape.table
                                break
                        except: continue
                
                if main_table:
                    start_row = 2
                    while len(main_table.rows) < (start_row + MAX_ROWS):
                        add_new_row(main_table)
                    
                    for idx_in_chunk, row_data in enumerate(chunk):
                        curr_idx = start_row + idx_in_chunk
                        cells = main_table.rows[curr_idx].cells
                        val_fluid = row_data.get(idx_fluid, '') if idx_fluid is not None else ''
                        _set_cell_text(cells[0], val_fluid, 9)
                        val_part = row_data.get(idx_part, '') if idx_part is not None else ''
                        _set_cell_text(cells[1], val_part, 9)
                        _set_cell_text(cells[2], "", 9) 
                        _set_cell_text(cells[3], row_data.get(idx_type, ''), 9)
                        _set_cell_text(cells[4], row_data.get(idx_spec, ''), 9)
                        _set_cell_text(cells[5], row_data.get(idx_grade, ''), 9)
                        val_ins = row_data.get(idx_insul, '') if idx_insul is not None else ''
                        _set_cell_text(cells[6], val_ins, 9)
                        val_temp = row_data.get(idx_temp, '') if idx_temp is not None else ''
                        _set_cell_text(cells[7], val_temp, 9)
                        val_press = row_data.get(idx_press, '') if idx_press is not None else ''
                        _set_cell_text(cells[8], val_press, 9)

                    for j in range(len(chunk), MAX_ROWS):
                        curr_idx = start_row + j
                        if curr_idx < len(main_table.rows):
                            for cell in main_table.rows[curr_idx].cells:
                                _set_cell_text(cell, "", 9)

            # --- 2. GENERATE DRAWING SLIDE FOR THIS GROUP ---
            # This happens immediately after the data slides for this specific file
            if current_drawing_filename:
                pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], current_drawing_filename)
                
                if os.path.exists(pdf_path):
                    try:
                        try: blank_slide_layout = prs.slide_layouts[6] 
                        except: blank_slide_layout = prs.slide_layouts[0] 
                            
                        slide_drawing = prs.slides.add_slide(blank_slide_layout)
                        
                        poppler_p = app.config.get('POPPLER_PATH') or os.environ.get('POPPLER_PATH')
                        images = pdf2image.convert_from_path(pdf_path, first_page=1, last_page=1, dpi=200, poppler_path=poppler_p)
                        
                        if images:
                            img = images[0]
                            width, height = img.size
                            
                            # === CROP LEFT 80% ===
                            left_crop = width * 0.02
                            top_crop = height * 0.02
                            right_crop = width * 0.82
                            bottom_crop = height * 0.98
                            
                            img_cropped = img.crop((left_crop, top_crop, right_crop, bottom_crop))
                            
                            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp_img:
                                img_cropped.save(tmp_img.name, 'JPEG')
                                tmp_img_path = tmp_img.name
                            
                            ppt_left = Inches(0.5)
                            ppt_top = Inches(0.5)
                            ppt_height = Inches(6.5) 
                            
                            slide_drawing.shapes.add_picture(tmp_img_path, ppt_left, ppt_top, height=ppt_height)
                            
                            try: os.remove(tmp_img_path)
                            except: pass

                    except Exception as e:
                        print(f"Error adding drawing slide: {e}")

        # --- PART 3: CLEANUP TEMPLATE ---
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[0])

        prs.save(output_pptx_path)
        return output_pptx_path
    
    except Exception as e:
        print(f"PPT Generation Error (Saving): {e}")
        return None
# ==================== MAIN ROUTE (UPDATED) ==================

@app.route('/generate_ppt/<filename>')
@login_required
def generate_ppt(filename):
    # This route now just calls the internal function
    path = generate_ppt_internal(filename)
    if path and os.path.exists(path):
        return send_file(path, as_attachment=True)
    else:
        flash("Error generating PowerPoint file.", "error")
        return redirect(url_for('index'))

# ==================== DATABASE MODELS ==================
class Role(db.Model):
    __tablename__ = 'roles'
    role_id = db.Column(db.Integer, primary_key=True)
    role_name = db.Column(db.String(50), unique=True, nullable=False)
    users = db.relationship('User', backref='role', lazy=True)

class User(db.Model, UserMixin):
    __tablename__ = 'users'
    user_id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    password_hash = db.Column(db.String(256), nullable=False)
    name = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    phone_num = db.Column(db.String(20), nullable=True)
    role_id = db.Column(db.Integer, db.ForeignKey('roles.role_id'), nullable=False)
    
    equipment_data = db.relationship('EquipmentData', backref='user', lazy=True)
    history_entries = db.relationship('History', backref='user', lazy=True)
    announcements = db.relationship('Announcement', backref='user', lazy=True)

    def get_id(self):
        return (self.user_id)
    def set_password(self, password):
        self.password_hash = generate_password_hash(password)
    def check_password(self, password):
        return check_password_hash(self.password_hash, password)
    def __repr__(self):
        return f'<User {self.username}>'

class History(db.Model):
    __tablename__ = 'history'
    history_id = db.Column(db.Integer, primary_key=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    excel_filename = db.Column(db.String(255), nullable=True)
    created_by_user_id = db.Column(db.Integer, db.ForeignKey('users.user_id'), nullable=False)
    equipment_data_rows = db.relationship('EquipmentData', backref='history_entry', lazy=True)

class EquipmentData(db.Model):
    __tablename__ = 'equipment_data'
    id = db.Column(db.Integer, primary_key=True)
    source_drawing = db.Column(db.String(255), nullable=False)
    part_name = db.Column(db.String(100), nullable=False)
    fluid = db.Column(db.String(100))
    material_type = db.Column(db.String(100))
    material_spec = db.Column(db.String(100))
    material_grade = db.Column(db.String(100))
    design_temp = db.Column(db.String(50))
    design_pressure = db.Column(db.String(50))
    operating_temp = db.Column(db.String(50))
    operating_pressure = db.Column(db.String(50))
    insulation = db.Column(db.String(50))
    no = db.Column(db.String(50), default='')
    equipment_no = db.Column(db.String(50), default='')
    pmt_no = db.Column(db.String(50), default='')
    equipment_description = db.Column(db.String(255), default='')
    phase = db.Column(db.String(50), default='')
    created_by_user_id = db.Column(db.Integer, db.ForeignKey('users.user_id'), nullable=True)
    history_id = db.Column(db.Integer, db.ForeignKey('history.history_id'), nullable=False)

class Announcement(db.Model):
    __tablename__ = 'announcements'
    announcement_id = db.Column(db.Integer, primary_key=True)
    message = db.Column(db.Text, nullable=False)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    attachment_filename = db.Column(db.String(255), nullable=True)
    user_id = db.Column(db.Integer, db.ForeignKey('users.user_id'), nullable=False)
    visible_to_manager = db.Column(db.Boolean, default=False, nullable=False)
    visible_to_engineer = db.Column(db.Boolean, default=False, nullable=False)

# ==================== FORMS ==================

class CreateUserForm(FlaskForm):
    username = StringField('Username', validators=[DataRequired(), Length(min=4, max=80)])
    role = SelectField('Role', choices=[('Admin', 'Admin'), ('Engineer', 'Engineer'), ('Manager', 'Manager')], validators=[DataRequired()])
    password = PasswordField('Password', validators=[DataRequired(), Length(min=6)])
    confirm_password = PasswordField('Confirm Password', validators=[DataRequired(), EqualTo('password')])
    submit = SubmitField('Create User')
    def validate_username(self, username):
        user = User.query.filter_by(username=username.data).first()
        if user: raise ValidationError('Username taken.')

class UpdateProfileForm(FlaskForm):
    username = StringField('Username', validators=[DataRequired(), Length(min=4, max=80)])
    name = StringField('Full Name', validators=[DataRequired(), Length(min=2, max=100)])
    email = StringField('Email', validators=[DataRequired(), Email()])
    phone_num = StringField('Phone Number', validators=[Optional(), Length(min=6, max=20)])
    submit = SubmitField('Save Changes')
    def validate_username(self, username):
        user = User.query.filter_by(username=username.data).first()
        if user and user.user_id != current_user.user_id: raise ValidationError('Username taken.')
    def validate_email(self, email):
        user = User.query.filter_by(email=email.data).first()
        if user and user.user_id != current_user.user_id: raise ValidationError('Email in use.')

class AnnouncementForm(FlaskForm):
    message = TextAreaField('Message', validators=[DataRequired(), Length(min=1, max=5000)])
    attachment = FileField('Attachment (Optional)', validators=[FileAllowed(['pdf', 'png', 'jpg', 'xlsx'])])
    visible_to_manager = BooleanField('Managers')
    visible_to_engineer = BooleanField('Engineers')
    submit = SubmitField('Post Announcement')
    def validate(self, **kwargs):
        if not super().validate(**kwargs): return False
        if not self.visible_to_manager.data and not self.visible_to_engineer.data:
            self.visible_to_manager.errors.append('Select at least one role.')
            return False
        return True

# ==================== HELPERS ==================

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role.role_name not in ['Manager', 'Admin']:
            flash('Permission denied.', 'error')
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

def manager_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role.role_name != 'Manager':
            flash('Permission denied to access this reporting page.', 'error')
            return redirect(url_for('index')) 
        return f(*args, **kwargs)
    return decorated_function

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

def get_gemini_prompt():
    return """
    You are a precise data extraction engine. Your task is to extract technical data and the full Bill of Materials (BOM) from the provided engineering drawing.

    1. GENERAL DATA:
       - Extract: Design Pressure, Design Temperature, Operating Pressure, Operating Temperature, Fluid, Insulation (Yes/No), Phase (Liquid/Gas).

    2. BILL OF MATERIALS (CRITICAL):
       - Locate the "Bill of Materials", "List of Materials", or "Nozzle Schedule" table.
       - Extract EVERY SINGLE ROW from this table exactly as written.
       - DO NOT FILTER OR SUMMARIZE. If the table has 20 rows, return 20 items.
       - For each row, map the columns to:
         - "part_name" (Description, Component, or Item Name)
         - "material_spec" (Material Specification, e.g., SA-516, SA-106)
         - "material_grade" (Grade, e.g., 70, B, 316L)

    Output strictly valid JSON matching this schema:
    {
      "design_pressure": "Value",
      "design_temperature": "Value",
      "operating_pressure": "Value",
      "operating_temperature": "Value",
      "fluid": "Value",
      "insulation": "Value",
      "phase": "Value",
      "parts_list": [
        { "part_name": "Shell", "material_spec": "SA-516", "material_grade": "70" },
        { "part_name": "Hex Nut", "material_spec": "SA-194", "material_grade": "2H" }
      ]
    }
    """

def call_gemini_api(images, prompt, api_key):
    api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
    headers = {"Content-Type": "application/json"}
    parts = [{"text": prompt}]
    for img in images:
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='JPEG')
        img_base64 = base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
        parts.append({"inline_data": {"mime_type": "image/jpeg", "data": img_base64}})
    
    # Updated Payload with Generation Config for JSON enforcement
    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "responseMimeType": "application/json"
        }
    }
    
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = requests.post(api_url, headers=headers, data=json.dumps(payload), timeout=125)
            response.raise_for_status()
            return response.json()['candidates'][0]['content']['parts'][0]['text']
        except Exception as e:
            if attempt == max_retries - 1: raise e
            sleep_time.sleep(1)

def clean_gemini_response(response_text):
    # Since we enforce JSON in call_gemini_api, this regex is a fallback
    match = re.search(r'```json\s*([\s\S]+?)\s*```', response_text)
    if match:
        return match.group(1)
    return response_text.strip()

def refine_material_type(spec, grade, ai_suggested_type="Not Found"):
    s = str(spec).upper() if spec else ""
    g = str(grade).upper() if grade else ""
    t = str(ai_suggested_type).upper() if ai_suggested_type else ""

    if "S275" in s or "S275" in g or "JR" in g: return "Structural Steel"
    if "193" in s or "193" in g: return "Stainless Steel Bolting"
    if "194" in s or "194" in g: return "Heavy Hex Nuts"
    if "G3507" in s or "G3507" in g: return "Carbon Steel"

    if "304" in g or "316" in g or "321" in g or "347" in g: return "Stainless Steel"
    if "SA-240" in s or "SA-312" in s or "SA-182" in s:
        if "F11" not in g and "F22" not in g: return "Stainless Steel"
    if "2205" in g or "S31803" in g: return "Duplex SS"
    if "SA-106" in s or "SA-105" in s or "SA-516" in s or "API 5L" in s or "A106" in s or "A516" in s: return "Carbon Steel"
    if "F11" in g or "F22" in g or "P11" in g or "P22" in g: return "Alloy Steel"

    if t and t != "NOT FOUND" and t != "NONE" and t != "OTHER": return ai_suggested_type

    return "Not Found"

def is_excluded_part(part_name):
    """
    Returns True if the part should be removed (Nuts, Bolts, Washers, etc.).
    This ensures consistent filtering in Python code, not AI guessing.
    """
    if not part_name or str(part_name).lower() == 'nan': 
        return True
        
    p = str(part_name).upper().strip()
    
    # LIST OF ITEMS TO IGNORE
    exclude_keywords = [
        'NUT', 'BOLT', 'WASHER', 'GASKET', 'STUD', 
        'NAME PLATE', 'NAMEPLATE', 'DATA PLATE',
        'SUPPORT', 'SADDLE', 'LEG', 'SKIRT', 
        'CLIP', 'BRACKET', 'LUG', 'LIFTING', 
        'EARTHING', 'BOSS', 'PAD', 'RIB', 'STIFFENER'
    ]
    
    # 1. Check Keywords
    if any(k in p for k in exclude_keywords):
        return True
        
    # 2. Check for vague/empty names
    if len(p) < 2 or p == "PLATE" or p == "PIPE":
        # Filter headers captured as data
        if p in ['ITEM', 'NO.', 'QTY', 'DESCRIPTION']: return True

    return False


def parse_gemini_response(json_text, drawing_name):
    extracted_data = []
    try: 
        data = json.loads(json_text)
    except: 
        data = {"parts_list": []}
    
    # Get the raw list (which now contains EVERYTHING, including nuts/bolts)
    parts = data.get("parts_list", []) or [{"part_name": "Not Found"}]
    
    fluid = data.get("fluid", "Not Found")
    phase = data.get("phase", "Not Found")
    insulation = data.get("insulation", "Not Found")
    design_temp = data.get("design_temperature", "Not Found")
    design_press = data.get("design_pressure", "Not Found")
    op_temp = data.get("operating_temperature", "Not Found")
    op_press = data.get("operating_pressure", "Not Found")

    for part in parts:
        raw_name = part.get("part_name", "Not Found")
        
        # --- NEW: DETERMINISTIC FILTERING ---
        # If it's a nut/bolt/gasket, SKIP IT.
        if is_excluded_part(raw_name):
            continue
        # ------------------------------------

        raw_spec = part.get("material_spec", "Not Found")
        raw_grade = part.get("material_grade", "Not Found")
        final_type = refine_material_type(raw_spec, raw_grade)

        extracted_data.append({
            "source_drawing": drawing_name,
            "part_name": raw_name,
            "fluid": fluid,
            "phase": phase,
            "insulation": insulation,
            "material_type": final_type, 
            "material_spec": raw_spec,
            "material_grade": raw_grade,
            'TEMP. (°C) ': design_temp,
            'PRESSURE (Mpa) ': design_press,
            'TEMP. (°C)': op_temp,
            'PRESSURE (Mpa)': op_press
        })
        
    return extracted_data

# ==================== ROUTES ==================

@app.route('/', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        if current_user.role.role_name == 'Admin':
            return redirect(url_for('admin_dashboard'))
        elif current_user.role.role_name == 'Manager':
            return redirect(url_for('manager_dashboard'))
        else:
            return redirect(url_for('index'))
    
    if request.method == 'POST':
        username_input = request.form.get('username')
        password_input = request.form.get('password')

        user = User.query.filter((User.username == username_input) | (User.email == username_input)).first()
        
        if user and user.check_password(password_input):
            login_user(user)
            if user.role.role_name == 'Admin':
                return redirect(url_for('admin_dashboard'))
            elif user.role.role_name == 'Manager':
                return redirect(url_for('manager_dashboard'))
            else:
                return redirect(url_for('index'))
            
        flash('Invalid username or password.', 'error')

    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

@app.route('/home')
@login_required
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
@login_required
def upload_file():
    if 'drawings' not in request.files:
        flash("No file part", "error")
        return redirect(url_for('index'))
    
    files = request.files.getlist('drawings')
    
    if not files or files[0].filename == '':
        flash("No selected file", "error")
        return redirect(url_for('index'))
    
    files.sort(key=lambda x: x.filename)
    all_data = [] 

    # --- DB CONNECTION FIX ---
    db.session.remove()
    
    try:
        for index, file in enumerate(files):
            if file.filename == '': continue
            
            safe_name = secure_filename(file.filename)
            path = os.path.join(app.config['UPLOAD_FOLDER'], safe_name)
            file.save(path)
            
            try:
                # Updated to use 300 DPI for better OCR/Vision accuracy
                images = pdf2image.convert_from_path(path, poppler_path=app.config.get('POPPLER_PATH') or POPPLER_PATH, dpi=300)
                text = call_gemini_api(images, get_gemini_prompt(), GEMINI_API_KEY)
                clean_text = clean_gemini_response(text)
                
                parsed = parse_gemini_response(clean_text, safe_name)
                
                if parsed:
                    all_data.extend(parsed) 
                
                if index < len(files) - 1:
                    sleep_time.sleep(20)

            except Exception as e:
                print(f"CRITICAL ERROR processing file {safe_name}: {e}")
                continue 
        
        if not all_data:
            flash("No data could be extracted. Check terminal for errors.", "error")
            return redirect(url_for('index'))
            
        temp_name = f"temp_{uuid.uuid4().hex}.json"
        with open(os.path.join(app.config['OUTPUT_FOLDER'], temp_name), 'w') as f:
            json.dump(all_data, f)
            
        return redirect(url_for('preview_data', temp_file=temp_name))

    except Exception as e:
        flash(f"System Error: {e}", "error")
        return redirect(url_for('index'))

@app.route('/preview')
@login_required
def preview_page():
    return render_template('preview.html', data_rows=[], equipment_count=0, confidence_score=0, missing_fields_count=0)

@app.route('/preview/<temp_file>')
@login_required
def preview_data(temp_file):
    try:
        with open(os.path.join(app.config['OUTPUT_FOLDER'], temp_file), 'r') as f:
            raw = json.load(f)
        
        preview_rows = []
        equip_count = 0
        curr_drawing = ""

        # --- STATISTICS CALCULATION VARIABLES ---
        total_critical_fields = 0
        filled_critical_fields = 0
        missing_fields_count = 0
        
        # Fields we expect the AI to find (Added 'insulation' and 'phase')
        critical_keys = [
            'part_name', 'fluid', 'material_type', 'material_spec', 
            'material_grade', 'TEMP. (°C) ', 'PRESSURE (Mpa) ',
            'insulation', 'phase'
        ]
        
        # List of values that define "Missing Data"
        missing_indicators = ['', 'NOT FOUND', 'NONE', '-', 'NAN', 'N/A', 'UNKNOWN', 'TBA', 'TBD', 'NOT APPLICABLE']

        for d in raw:
            if d['source_drawing'] != curr_drawing:
                equip_count += 1
                curr_drawing = d['source_drawing']
            
            # --- CALCULATE SCORE FOR THIS ROW ---
            for k in critical_keys:
                total_critical_fields += 1
                # Clean the value: remove dots, extra spaces, convert to upper
                val = str(d.get(k, '')).strip().upper()
                
                # Broadened Check:
                # 1. Exact match in list
                # 2. Substring match (e.g. "PART NOT FOUND" contains "NOT FOUND")
                is_missing = False
                if val in missing_indicators:
                    is_missing = True
                elif 'NOT FOUND' in val or 'UNKNOWN' in val:
                    is_missing = True
                
                if is_missing:
                    missing_fields_count += 1
                else:
                    filled_critical_fields += 1
            # ------------------------------------
            
            row = {
                'NO.': equip_count,
                'EQUIPMENT NO. ': f"V-{equip_count:03d}",
                'PMT NO.': os.path.splitext(d['source_drawing'])[0],
                'EQUIPMENT DESCRIPTION': "",
                'PARTS': d.get('part_name'),
                'PHASE': d.get('phase', 'Not Found'),
                'FLUID': d.get('fluid'),
                'TYPE': d.get('material_type'),
                'SPEC.': d.get('material_spec'),
                'GRADE': d.get('material_grade'),
                'INSULATION': d.get('insulation', 'Not Found'),
                'TEMP. (°C) ': d.get('TEMP. (°C) '),
                'PRESSURE (Mpa) ': d.get('PRESSURE (Mpa) '),
                'TEMP. (°C)': d.get('TEMP. (°C)'),
                'PRESSURE (Mpa)': d.get('PRESSURE (Mpa)'),
                'source_drawing': d.get('source_drawing')
            }
            preview_rows.append(row)
            
        # --- FINAL CALCULATION ---
        if total_critical_fields > 0:
            confidence_score = round((filled_critical_fields / total_critical_fields) * 100, 1)
        else:
            confidence_score = 0
        # -------------------------

        # --- NEW: GET UNIQUE SOURCE FILES FOR THE PDF VIEWER ---
        # Extract unique source_drawing names using a set, then convert back to list
        unique_files = list(set(row['source_drawing'] for row in preview_rows if row.get('source_drawing')))
        unique_files.sort() # Sort them alphabetically
        # -------------------------------------------------------
            
        return render_template(
            'preview.html', 
            data_rows=preview_rows, 
            equipment_count=equip_count, 
            temp_file=temp_file,
            confidence_score=confidence_score,
            missing_fields_count=missing_fields_count,
            unique_files=unique_files # <--- Pass this new variable to the template
        )
    except Exception as e:
        flash(f"Error loading preview: {e}", "error")
        return redirect(url_for('index'))

@app.route('/view_uploaded_file/<filename>')
@login_required
def view_uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], secure_filename(filename), as_attachment=False)

@app.route('/manual-input')
@login_required
def manual_input():
    return render_template('manual_input.html')

@app.route('/save_data', methods=['POST'])
@login_required
def save_data():
    try:
        def get_vals(name): return [x.strip() for x in request.form.getlist(name)]
        
        rows = []
        parts = get_vals('PARTS') 
        
        for i in range(len(parts)):
            rows.append({
                'NO.': get_vals('NO.')[i],
                'EQUIPMENT NO. ': get_vals('EQUIPMENT NO. ')[i],
                'PMT NO.': get_vals('PMT NO.')[i],
                'EQUIPMENT DESCRIPTION': get_vals('EQUIPMENT DESCRIPTION')[i],
                'PARTS': parts[i],
                'PHASE': get_vals('PHASE')[i],
                'FLUID': get_vals('FLUID')[i],
                'TYPE': get_vals('TYPE')[i],
                'SPEC.': get_vals('SPEC.')[i],
                'GRADE': get_vals('GRADE')[i],
                'INSULATION': get_vals('INSULATION')[i],
                'TEMP. (°C) ': get_vals('TEMP. (°C) ')[i],
                'PRESSURE (Mpa) ': get_vals('PRESSURE (Mpa) ')[i],
                'TEMP. (°C)': get_vals('TEMP. (°C)')[i],
                'PRESSURE (Mpa)': get_vals('PRESSURE (Mpa)')[i],
                'source_drawing': request.form.getlist('source_drawing')[i]
            })

        history = History(created_by_user_id=current_user.user_id)
        db.session.add(history)
        db.session.flush()
        
        for r in rows:
            db.session.add(EquipmentData(
                history_id=history.history_id,
                created_by_user_id=current_user.user_id,
                source_drawing=r['source_drawing'],
                part_name=r['PARTS'],
                fluid=r['FLUID'],
                material_type=r['TYPE'],
                material_spec=r['SPEC.'],
                material_grade=r['GRADE'],
                design_temp=r['TEMP. (°C) '],
                design_pressure=r['PRESSURE (Mpa) '],
                operating_temp=r['TEMP. (°C)'],
                operating_pressure=r['PRESSURE (Mpa)'],
                insulation=r['INSULATION'],
                no=r['NO.'],
                equipment_no=r['EQUIPMENT NO. '],
                pmt_no=r['PMT NO.'],
                equipment_description=r['EQUIPMENT DESCRIPTION'],
                phase=r['PHASE']
            ))

        excel_rows = []
        prev_equip_no = None
        for r in rows:
            row_copy = r.copy()
            curr_equip = r['EQUIPMENT NO. ']
            if curr_equip == prev_equip_no:
                row_copy['NO.'] = ""
                row_copy['EQUIPMENT NO. '] = ""
                row_copy['PMT NO.'] = ""
                row_copy['EQUIPMENT DESCRIPTION'] = ""
                row_copy['PHASE'] = ""
            else:
                prev_equip_no = curr_equip
            excel_rows.append(row_copy)

        # --- NEW NAMING LOGIC START ---
        # Get the name of the first drawing in this batch to use as the filename
        first_drawing_name = rows[0]['source_drawing'] 
        # Remove the .pdf extension and clean it up
        clean_name = os.path.splitext(first_drawing_name)[0]
        # Create a meaningful filename (e.g., MLK_PMT_10103... .xlsx)
        filename = f"{clean_name}.xlsx"
        filepath = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        # --- NEW NAMING LOGIC END ---
        
        shutil.copyfile(TEMPLATE_FILE, filepath)
        book = load_workbook(filepath)
        start_row = book[TARGET_SHEET].max_row
        book.close()
        
        df = pd.DataFrame(excel_rows).reindex(columns=template_column_order)
        
        with pd.ExcelWriter(filepath, engine='openpyxl', mode='a', if_sheet_exists='overlay') as writer:
            df.to_excel(writer, sheet_name=TARGET_SHEET, startrow=start_row, index=False, header=False)

        wb = load_workbook(filepath)
        ws = wb[TARGET_SHEET]
        cols_to_merge = [1, 2, 3, 4]
        row = 2 
        max_row = ws.max_row

        while row <= max_row:
            cell_val = ws.cell(row=row, column=2).value
            if cell_val:
                start_merge_row = row
                next_row = row + 1
                while next_row <= max_row and not ws.cell(row=next_row, column=2).value:
                    next_row += 1
                end_merge_row = next_row - 1
                if end_merge_row > start_merge_row:
                    for col_idx in cols_to_merge:
                        ws.merge_cells(start_row=start_merge_row, start_column=col_idx, end_row=end_merge_row, end_column=col_idx)
                        cell = ws.cell(row=start_merge_row, column=col_idx)
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                else:
                    for col_idx in cols_to_merge:
                        cell = ws.cell(row=start_merge_row, column=col_idx)
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                row = next_row
            else:
                row += 1
        
        wb.save(filepath)

        history.excel_filename = filename
        db.session.commit()
        
        if request.form.get('temp_file'):
            try: os.remove(os.path.join(app.config['OUTPUT_FOLDER'], request.form.get('temp_file')))
            except: pass
            
        flash('Data saved successfully!', 'success')
        return render_template('preview.html', data_rows=rows, equipment_count=len(set(get_vals('EQUIPMENT NO. '))), excel_file=filename, confidence_score=0, missing_fields_count=0)

    except Exception as e:
        db.session.rollback()
        flash(f"Error saving: {e}", "error")
        return redirect(url_for('manual_input'))

@app.route('/download/<filename>')
@login_required
def download_file(filename):
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename, as_attachment=True)

@app.route('/notifications')
@login_required
def notifications():
    anns = Announcement.query.filter((Announcement.visible_to_engineer == True) | (Announcement.visible_to_manager == True)).order_by(Announcement.created_at.desc()).all()
    return render_template('notification.html', announcements=anns)

@app.route('/profile', methods=['GET', 'POST'])
@login_required
def personal_info():
    form = UpdateProfileForm()
    if form.validate_on_submit():
        current_user.username = form.username.data
        current_user.name = form.name.data
        current_user.email = form.email.data
        current_user.phone_num = form.phone_num.data
        db.session.commit()
        flash('Profile updated.', 'success')
        return redirect(url_for('personal_info'))
    form.username.data = current_user.username
    form.name.data = current_user.name
    form.email.data = current_user.email
    form.phone_num.data = current_user.phone_num
    return render_template('personal_info.html', form=form)

@app.route('/download/announcement/<filename>')
@login_required
def download_announcement(filename):
    return send_from_directory(app.config['ANNOUNCEMENT_FOLDER'], filename, as_attachment=True)

@app.route('/send_files_email', methods=['POST'])
@login_required
def send_files_email():
    recipient = request.form.get('recipient_email')
    excel_file = request.form.get('excel_filename')
    
    excel_path = os.path.join(app.config['OUTPUT_FOLDER'], excel_file)
    ppt_filename = f"Inspection_Plan_{excel_file.replace('.xlsx', '.pptx')}"
    ppt_path = os.path.join(app.config['OUTPUT_FOLDER'], ppt_filename)

    if not recipient:
        flash("Please enter an email address.", "error")
        return redirect(url_for('index'))

    if not os.path.exists(ppt_path):
        print("PPT missing, generating now...")
        generated_path = generate_ppt_internal(excel_file)
        if not generated_path:
            flash("Warning: Could not generate PowerPoint attachment.", "error")

    try:
        msg = Message("Generated Inspection Files", recipients=[recipient])
        msg.body = f"Hello,\n\nPlease find attached the generated Inspection Plan and Excel Data for {excel_file}.\n\nSent from iPetro System."

        if os.path.exists(excel_path):
            with app.open_resource(excel_path) as fp:
                msg.attach(excel_file, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", fp.read())
        
        if os.path.exists(ppt_path):
            with app.open_resource(ppt_path) as fp:
                msg.attach(ppt_filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation", fp.read())
        else:
            msg.body += "\n\n(Note: PowerPoint file could not be generated and is missing from this email.)"

        mail.send(msg)
        flash(f"Email successfully sent to {recipient}!", "success")
        
    except Exception as e:
        print(f"Email Error: {e}")
        flash(f"Failed to send email: {e}", "error")

    return redirect(url_for('index'))
# ==================== MANAGER ROUTES ==================

@app.route('/manager/dashboard')
@login_required
@manager_required
def manager_dashboard():
    total_batches = History.query.count()
    total_parts = EquipmentData.query.count()
    engineer_role = Role.query.filter_by(role_name='Engineer').first()
    engineer_count = User.query.filter_by(role_id=engineer_role.role_id).count() if engineer_role else 0
    recent_activity = History.query.order_by(History.created_at.desc()).limit(5).all()
    
    top_engineers = db.session.query(
        User.name, 
        func.count(History.history_id).label('batch_count')
    ).join(History, User.user_id == History.created_by_user_id
    ).filter(User.role_id == engineer_role.role_id if engineer_role else True).group_by(User.name
    ).order_by(func.count(History.history_id).desc()
    ).limit(3).all()

    return render_template(
        'dashboard_manager.html',
        total_batches=total_batches,
        total_parts=total_parts,
        engineer_count=engineer_count,
        recent_activity=recent_activity,
        top_engineers=top_engineers
    )

@app.route('/manager/history')
@login_required
@manager_required
def manager_history():
    return render_template('manager_history.html', all_history=History.query.order_by(History.created_at.desc()).all())

@app.route('/manager/reports')
@login_required
@manager_required
def manager_reports():
    # --- CHANGED: Count distinct source drawings instead of equipment numbers ---
    total_files = db.session.query(func.count(func.distinct(EquipmentData.source_drawing))).scalar() or 0

    total_components = EquipmentData.query.count() or 0
    
    bad_data_count = EquipmentData.query.filter(
        (EquipmentData.material_type.ilike('%Not Found%')) | 
        (EquipmentData.material_spec.ilike('%Not Found%'))
    ).count()
    
    success_rate = 0
    if total_components > 0:
        success_rate = round(((total_components - bad_data_count) / total_components) * 100, 1)

    material_stats = db.session.query(
        EquipmentData.material_type, 
        func.count(EquipmentData.id)
    ).group_by(EquipmentData.material_type).all()
    
    mat_labels = []
    mat_counts = []
    for m in material_stats:
        label = m[0] if m[0] and m[0].strip() != '' else 'Unknown'
        mat_labels.append(label)
        mat_counts.append(m[1])

    action_items = EquipmentData.query.filter(
        (EquipmentData.material_type.ilike('%Not Found%')) | 
        (EquipmentData.material_spec.ilike('%Not Found%'))
    ).limit(50).all()

    fluids = [r[0] for r in db.session.query(EquipmentData.fluid).distinct().all() if r[0]]

    equipment_list = db.session.query(
        EquipmentData.equipment_no,
        func.max(EquipmentData.equipment_description).label('description'),
        func.max(EquipmentData.pmt_no).label('pmt')
    ).filter(EquipmentData.equipment_no != '').group_by(EquipmentData.equipment_no).all()

    return render_template(
        'report.html',
        total_files=total_files,
        total_components=total_components,
        success_rate=success_rate,
        mat_labels=json.dumps(mat_labels),
        mat_counts=json.dumps(mat_counts),
        action_items=action_items,
        fluids=fluids,
        equipment_list=equipment_list,
        current_date=datetime.now().strftime("%d %B %Y")
    )

@app.route('/manager/review')
@login_required
@manager_required
def manager_review_queue():
    return render_template('queue.html')

# ==================== ADMIN ROUTES ==================

@app.route('/admin/dashboard')
@login_required
@admin_required
def admin_dashboard():
    u_count = User.query.count()
    admin_role = Role.query.filter_by(role_name='Admin').first()
    admin_count = User.query.filter_by(role_id=admin_role.role_id).count() if admin_role else 0
    manager_role = Role.query.filter_by(role_name='Manager').first()
    manager_count = User.query.filter_by(role_id=manager_role.role_id).count() if manager_role else 0
    engineer_role = Role.query.filter_by(role_name='Engineer').first()
    engineer_count = User.query.filter_by(role_id=engineer_role.role_id).count() if engineer_role else 0
    f_total = History.query.count()
    today = datetime.utcnow().date()
    f_today = History.query.filter(History.created_at >= datetime.combine(today, time.min), History.created_at <= datetime.combine(today, time.max)).count()
    recent = User.query.order_by(User.user_id.desc()).limit(5).all()
    return render_template('dashboard_admin.html', user_count=u_count, admin_count=admin_count, manager_count=manager_count, engineer_count=engineer_count, file_count_total=f_total, file_count_today=f_today, recent_users=recent)

@app.route('/admin/create-user', methods=['GET', 'POST'])
@login_required
@admin_required
def admin_create_user():
    form = CreateUserForm()
    if form.validate_on_submit():
        role = Role.query.filter_by(role_name=form.role.data).first()
        u = User(username=form.username.data, name=form.username.data, email=f"{form.username.data}@ipetro.com", role_id=role.role_id)
        u.set_password(form.password.data)
        db.session.add(u)
        db.session.commit()
        flash('User created.', 'success')
        return redirect(url_for('admin_create_user'))
    return render_template('createuser.html', form=form)

@app.route('/admin/announcement', methods=['GET', 'POST'])
@login_required
@admin_required
def admin_announcement():
    form = AnnouncementForm()
    if form.validate_on_submit():
        fname = None
        if form.attachment.data:
            fname = secure_filename(form.attachment.data.filename)
            form.attachment.data.save(os.path.join(app.config['ANNOUNCEMENT_FOLDER'], fname))
        db.session.add(Announcement(message=form.message.data, attachment_filename=fname, user_id=current_user.user_id, visible_to_manager=form.visible_to_manager.data, visible_to_engineer=form.visible_to_engineer.data))
        db.session.commit()
        flash('Announcement posted.', 'success')
        return redirect(url_for('admin_announcement'))
    return render_template('announcement.html', form=form, announcements=Announcement.query.order_by(Announcement.created_at.desc()).all())

@app.route('/admin/history')
@login_required
@admin_required
def admin_history():
    items = History.query.order_by(History.created_at.desc()).all()
    return render_template('history_admin.html', histories=items)

@app.route('/admin/reports')
@login_required
@admin_required
def admin_reports():
    # --- CHANGED: Count distinct source drawings instead of equipment numbers ---
    total_files = db.session.query(func.count(func.distinct(EquipmentData.source_drawing))).scalar() or 0
    
    total_components = EquipmentData.query.count() or 0
    
    bad_data_count = EquipmentData.query.filter(
        (EquipmentData.material_type.ilike('%Not Found%')) | 
        (EquipmentData.material_spec.ilike('%Not Found%'))
    ).count()
    
    success_rate = 0
    if total_components > 0:
        success_rate = round(((total_components - bad_data_count) / total_components) * 100, 1)

    material_stats = db.session.query(
        EquipmentData.material_type, 
        func.count(EquipmentData.id)
    ).group_by(EquipmentData.material_type).all()
    
    mat_labels = []
    mat_counts = []
    for m in material_stats:
        label = m[0] if m[0] and m[0].strip() != '' else 'Unknown'
        mat_labels.append(label)
        mat_counts.append(m[1])

    action_items = EquipmentData.query.filter(
        (EquipmentData.material_type.ilike('%Not Found%')) | 
        (EquipmentData.material_spec.ilike('%Not Found%'))
    ).limit(50).all()

    fluids = [r[0] for r in db.session.query(EquipmentData.fluid).distinct().all() if r[0]]

    equipment_list = db.session.query(
        EquipmentData.equipment_no,
        func.max(EquipmentData.equipment_description).label('description'),
        func.max(EquipmentData.pmt_no).label('pmt')
    ).filter(EquipmentData.equipment_no != '').group_by(EquipmentData.equipment_no).all()

    return render_template(
        'report.html',
        total_files=total_files,
        total_components=total_components,
        success_rate=success_rate,
        mat_labels=json.dumps(mat_labels),
        mat_counts=json.dumps(mat_counts),
        action_items=action_items,
        fluids=fluids,
        equipment_list=equipment_list,
        current_date=datetime.now().strftime("%d %B %Y")
    )

@app.route('/admin/review-queue')
@login_required
@admin_required
def admin_review_queue():
    return render_template('queue.html')

@app.route('/admin/statistics')
@login_required
@admin_required
def admin_statistics():
    flash("The Statistics page is currently under construction.", "info")
    return redirect(url_for('admin_dashboard'))

@app.cli.command('init-db')
@click.option('--drop', is_flag=True)
def init_db(drop):
    if drop: db.drop_all()
    db.create_all()
    if not Role.query.first():
        m = Role(role_name='Manager')
        e = Role(role_name='Engineer')
        a = Role(role_name="Admin")
        db.session.add_all([m, e, a])
        db.session.commit()
        
        manager_user = User(username='manager@ipetro.com', name='Manager', email='manager@ipetro.com', role_id=m.role_id)
        manager_user.set_password('abc1234')
        
        eng_user = User(username='engineer@ipetro.com', name='Eng', email='engineer@ipetro.com', role_id=e.role_id)
        eng_user.set_password('abc1234')
        
        admin_user = User(username='admin@ipetro.com', name='Admin', email='Admin@ipetro.com', role_id=a.role_id)
        admin_user.set_password('abc1234')
        
        db.session.add_all([manager_user, eng_user, admin_user])
        db.session.commit()
    print("DB Initialized.")

if __name__ == '__main__':
    app.run(debug=True)