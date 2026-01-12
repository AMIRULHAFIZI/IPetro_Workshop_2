import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

if len(sys.argv) < 2:
    print("Usage: inspect_pptx.py <path-to-pptx>")
    raise SystemExit(1)

INPUT = Path(sys.argv[1])
if not INPUT.exists():
    print(f"Input not found: {INPUT}")
    raise SystemExit(1)

prs = Presentation(str(INPUT))
for i, slide in enumerate(prs.slides):
    print(f"Slide {i} (shapes: {len(slide.shapes)})")
    for j, shape in enumerate(slide.shapes):
        try:
            stype = shape.shape_type
        except Exception:
            stype = 'unknown'
        try:
            has_image = getattr(shape, 'image', None) is not None
        except Exception:
            has_image = False
        try:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
        except Exception:
            left = top = width = height = 'n/a'
        name = getattr(shape, 'name', '')
        print(f"  Shape {j}: type={stype}, has_image={has_image}, name='{name}', left={left}, top={top}, w={width}, h={height}")
    print('')
print('Done')
