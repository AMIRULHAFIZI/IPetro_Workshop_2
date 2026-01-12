import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from pathlib import Path

if len(sys.argv) < 2:
    print("Usage: move_drawing.py <input-pptx> [output-pptx]")
    raise SystemExit(1)

INPUT = Path(sys.argv[1])
if len(sys.argv) >= 3:
    OUTPUT = Path(sys.argv[2])
else:
    OUTPUT = INPUT.with_name(INPUT.stem + "_drawing_slide.pptx")

if not INPUT.exists():
    print(f"Input file not found: {INPUT}")
    raise SystemExit(1)

prs = Presentation(str(INPUT))
slide_width = prs.slide_width
slide_height = prs.slide_height

found = None
found_slide_idx = None
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # prefer pictures on the right half of the slide (likely drawing)
                if shape.left > slide_width // 2:
                    found = shape
                    found_slide_idx = i
                    break
                if found is None:
                    found = shape
                    found_slide_idx = i
        except Exception:
            try:
                _ = getattr(shape, 'image', None)
                if _ is not None and found is None:
                    found = shape
                    found_slide_idx = i
            except Exception:
                pass
    if found is not None and found_slide_idx == i:
        break

if found is None:
    print("No picture shape found in presentation.")
    raise SystemExit(1)

# extract image bytes
try:
    img_blob = found.image.blob
except Exception:
    print("Found shape did not expose image blob.")
    raise SystemExit(1)

# Append a new blank slide to a copy of the original presentation and add the picture
# (this preserves the original file and saves a new file with the extra drawing slide)
# We'll add the slide to the loaded `prs` and save to OUTPUT.
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
new_slide = prs.slides.add_slide(blank_layout)

# add picture at full slide width, then vertically center it
img_stream = BytesIO(img_blob)
pic = new_slide.shapes.add_picture(img_stream, 0, 0, width=slide_width)
try:
    # center vertically
    pic.top = int((slide_height - pic.height) // 2)
except Exception:
    pass

prs.save(str(OUTPUT))
print(f"Saved new presentation (original + drawing slide): {OUTPUT}")
print(f"Original slide index containing drawing: {found_slide_idx}")
