from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

root = Path('.')
for p in root.rglob('*.pptx'):
    try:
        prs = Presentation(str(p))
    except Exception as e:
        print(f"{p}: failed to open ({e})")
        continue
    found_types = set()
    picture_count = 0
    group_count = 0
    ole_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                stype = shape.shape_type
            except Exception:
                stype = None
            found_types.add(str(stype))
            if stype is not None:
                if stype.name == 'PICTURE':
                    picture_count += 1
                if stype.name == 'GROUP':
                    group_count += 1
                if stype.name == 'OLE_OBJECT':
                    ole_count += 1
    print(f"{p}: pictures={picture_count}, groups={group_count}, ole={ole_count}, types={sorted(list(found_types))}")
